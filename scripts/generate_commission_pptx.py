#!/usr/bin/env python3
"""
NEBULA Agency — Générateur PPTX Présentation Conférence
Structure des Commissions 2025 — Format 16:9 plein écran
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu, Cm
from pptx.enum.dml import MSO_THEME_COLOR
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# ─── Dimensions 16:9 ──────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

# ─── Palette NEBULA ───────────────────────────────────────────────
C = {
    'dark':    RGBColor(0x0A, 0x0A, 0x1A),
    'card':    RGBColor(0x1E, 0x1B, 0x4B),
    'card2':   RGBColor(0x13, 0x11, 0x1F),
    'purple':  RGBColor(0x6C, 0x3C, 0xE1),
    'violet':  RGBColor(0x8B, 0x5C, 0xF6),
    'pink':    RGBColor(0xEC, 0x48, 0x99),
    'cyan':    RGBColor(0x06, 0xB6, 0xD4),
    'gold':    RGBColor(0xF5, 0x9E, 0x0B),
    'silver':  RGBColor(0x94, 0xA3, 0xB8),
    'green':   RGBColor(0x10, 0xB9, 0x81),
    'white':   RGBColor(0xFF, 0xFF, 0xFF),
    'muted':   RGBColor(0x64, 0x74, 0x8B),
    'starter': RGBColor(0x3B, 0x82, 0xF6),
    'border':  RGBColor(0x37, 0x30, 0xA3),
    'light':   RGBColor(0xF8, 0xFA, 0xFC),
}


# ─── Helpers ──────────────────────────────────────────────────────

def rgb(r, g, b):
    return RGBColor(r, g, b)


def add_rect(slide, x, y, w, h, fill_color, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape


def add_rounded_rect(slide, x, y, w, h, fill_color, radius_pt=8, line_color=None, line_width=0):
    shape = slide.shapes.add_shape(5, x, y, w, h)  # 5 = rounded rectangle
    w_in = w / 914400 if isinstance(w, (int, float)) else w.inches
    h_in = h / 914400 if isinstance(h, (int, float)) else h.inches
    ref = min(w_in * 96, h_in * 96)
    adj = int(radius_pt * 914.4 / ref * 100000) if ref > 0 else 0
    adj = max(0, min(adj, 50000))
    shape.adjustments[0] = adj / 100000
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h, size=18, bold=False, color=None,
             align=PP_ALIGN.LEFT, font='Calibri', italic=False):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    if color:
        run.font.color.rgb = color
    return txBox


def add_multiline_text(slide, lines, x, y, w, h, default_size=18,
                       default_color=None, default_bold=False, align=PP_ALIGN.LEFT):
    """lines = list of (text, size, bold, color) tuples"""
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for line_data in lines:
        if isinstance(line_data, str):
            text, size, bold, color = line_data, default_size, default_bold, default_color
        else:
            text = line_data[0]
            size = line_data[1] if len(line_data) > 1 else default_size
            bold = line_data[2] if len(line_data) > 2 else default_bold
            color = line_data[3] if len(line_data) > 3 else default_color

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = 'Calibri'
        if color:
            run.font.color.rgb = color
    return txBox


def bg_dark(slide):
    """Fond sombre dégradé simulé"""
    add_rect(slide, 0, 0, W, H, C['dark'])
    # Accent violet en haut à gauche
    add_rect(slide, 0, 0, Inches(4), Inches(0.06), C['purple'])


def add_badge(slide, text, x, y, bg_color, text_color=None):
    if text_color is None:
        text_color = C['white']
    w = Inches(2.2)
    h = Inches(0.32)
    add_rounded_rect(slide, x, y, w, h, bg_color, radius_pt=12)
    add_text(slide, text, x, y + Inches(0.04), w, h - Inches(0.04),
             size=10, bold=True, color=text_color, align=PP_ALIGN.CENTER)


def add_pill(slide, text, x, y, w, h, bg_color, txt_color, size=14):
    add_rounded_rect(slide, x, y, w, h, bg_color, radius_pt=14)
    add_text(slide, text, x, y, w, h, size=size, bold=True,
             color=txt_color, align=PP_ALIGN.CENTER)


def footer(slide, page_num, total, label=''):
    # Ligne bas
    add_rect(slide, 0, H - Inches(0.04), W, Inches(0.04), C['purple'])
    # Logo bas gauche
    add_rounded_rect(slide, Inches(0.3), H - Inches(0.48), Inches(2.0), Inches(0.32),
                     C['purple'], radius_pt=8)
    add_text(slide, 'NEBULA AGENCY', Inches(0.3), H - Inches(0.48), Inches(2.0), Inches(0.32),
             size=9, bold=True, color=C['white'], align=PP_ALIGN.CENTER)
    # Numéro slide
    add_text(slide, f'{page_num} / {total}', W - Inches(1.2), H - Inches(0.48),
             Inches(1.0), Inches(0.32), size=10, color=C['muted'], align=PP_ALIGN.RIGHT)
    if label:
        add_text(slide, label, Inches(2.6), H - Inches(0.48),
                 W - Inches(4.0), Inches(0.32), size=9, color=C['muted'], align=PP_ALIGN.CENTER)


def add_circle_node(slide, cx, cy, r, fill, label_top, label_bot, size_top=13, size_bot=10,
                    initials='', highlight=False, badge_text=''):
    """Cercle avec label haut+bas — cx/cy en Inches"""
    ri = Inches(r)
    xi = Inches(cx) - ri
    yi = Inches(cy) - ri

    if highlight:
        # Halo
        add_rounded_rect(slide, xi - Inches(0.08), yi - Inches(0.08),
                         ri * 2 + Inches(0.16), ri * 2 + Inches(0.16),
                         RGBColor(0xF5, 0x9E, 0x0B), radius_pt=999)

    s = slide.shapes.add_shape(9, xi, yi, ri * 2, ri * 2)  # 9 = ellipse
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()

    if initials:
        add_text(slide, initials, xi, yi + ri - Pt(10), ri * 2, ri,
                 size=int(r * 22), bold=True, color=C['white'], align=PP_ALIGN.CENTER)

    add_text(slide, label_top,
             Inches(cx) - Inches(1.0), Inches(cy) + ri + Inches(0.04),
             Inches(2.0), Inches(0.28),
             size=size_top, bold=True, color=C['white'], align=PP_ALIGN.CENTER)

    if label_bot:
        add_text(slide, label_bot,
                 Inches(cx) - Inches(1.0), Inches(cy) + ri + Inches(0.30),
                 Inches(2.0), Inches(0.24),
                 size=size_bot, bold=False, color=fill, align=PP_ALIGN.CENTER)

    if badge_text:
        bw = Inches(0.6)
        bh = Inches(0.24)
        add_rounded_rect(slide, Inches(cx) + ri - Inches(0.1), Inches(cy) - ri - Inches(0.1),
                         bw, bh, C['gold'], radius_pt=8)
        add_text(slide, badge_text,
                 Inches(cx) + ri - Inches(0.1), Inches(cy) - ri - Inches(0.1),
                 bw, bh, size=8, bold=True, color=C['dark'], align=PP_ALIGN.CENTER)


def add_connector_line(slide, x1, y1, x2, y2, color, width_pt=2):
    """Ligne connecteur entre deux points (en Inches)"""
    # Utilise add_shape type line
    connector = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    connector.line.color.rgb = color
    connector.line.width = Pt(width_pt)


def add_label_pill_on_line(slide, cx, cy, text, color):
    w, h = Inches(1.1), Inches(0.26)
    add_rounded_rect(slide, Inches(cx) - w/2, Inches(cy) - h/2, w, h,
                     C['dark'], radius_pt=8)
    add_text(slide, text, Inches(cx) - w/2, Inches(cy) - h/2 + Inches(0.01),
             w, h, size=8, bold=True, color=color, align=PP_ALIGN.CENTER)


# ─── SLIDES ───────────────────────────────────────────────────────

def slide_01_cover(prs):
    """Slide 1 — Couverture"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg_dark(slide)

    # Grand cercle décoratif fond
    s = slide.shapes.add_shape(9, W - Inches(4.5), Inches(-1.0), Inches(6.0), Inches(6.0))
    s.fill.solid(); s.fill.fore_color.rgb = RGBColor(0x2D, 0x1A, 0x6E)
    s.line.fill.background()

    s2 = slide.shapes.add_shape(9, Inches(-1.0), H - Inches(3.0), Inches(3.5), Inches(3.5))
    s2.fill.solid(); s2.fill.fore_color.rgb = RGBColor(0x1A, 0x0D, 0x45)
    s2.line.fill.background()

    # Badge
    add_badge(slide, 'NEBULA AGENCY — PARTENARIAT 2025',
              Inches(0.8), Inches(1.2), C['purple'])

    # Titre principal
    add_text(slide, 'PROGRAMME DE', Inches(0.8), Inches(1.8),
             Inches(8.0), Inches(0.9), size=48, bold=True, color=C['white'])
    add_text(slide, 'COMMISSION', Inches(0.8), Inches(2.6),
             Inches(8.5), Inches(1.1), size=64, bold=True, color=C['violet'])

    # Sous-titre
    add_text(slide, 'Vendez des vitrines digitales & construisez un réseau qui travaille pour vous',
             Inches(0.8), Inches(3.8), Inches(7.5), Inches(0.8),
             size=18, color=C['silver'])

    # 3 pills stats
    pill_data = [
        ('25% → 35%', 'Commission directe', C['gold'], Inches(1.0)),
        ('3 niveaux', 'Revenus de réseau', C['violet'], Inches(4.0)),
        ('48h', 'Délai de paiement', C['green'], Inches(7.0)),
    ]
    for val, label, col, px in pill_data:
        add_rounded_rect(slide, px, Inches(5.0), Inches(2.4), Inches(1.6),
                         C['card'], radius_pt=12)
        add_rect(slide, px, Inches(5.0), Inches(2.4), Inches(0.05), col)
        add_text(slide, val, px, Inches(5.1), Inches(2.4), Inches(0.7),
                 size=26, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_text(slide, label, px, Inches(5.75), Inches(2.4), Inches(0.5),
                 size=11, color=C['silver'], align=PP_ALIGN.CENTER)

    footer(slide, 1, 12, 'Structure des commissions — Confidentiel partenaires')


def slide_02_agenda(prs):
    """Slide 2 — Plan de la présentation"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_dark(slide)

    add_text(slide, 'AU PROGRAMME', Inches(0.7), Inches(0.4),
             Inches(8), Inches(0.6), size=13, bold=True, color=C['violet'])
    add_text(slide, 'Ce que vous allez découvrir',
             Inches(0.7), Inches(0.85), Inches(10), Inches(0.8),
             size=36, bold=True, color=C['white'])

    items = [
        ('01', 'Les 3 paliers de commission', 'STARTER, SILVER, GOLD — comment progresser', C['starter']),
        ('02', 'Le réseau 3 niveaux', 'N1, N2, N3 — gagner sans vendre', C['violet']),
        ('03', 'Exemple concret — Shad', "L'arbre réseau et les gains réels du mois", C['gold']),
        ('04', 'Ce que Marc gagne en passif', 'Le pouvoir des revenus résiduels', C['pink']),
        ('05', 'Simulation de vos revenus', 'De débutant à GOLD Élite', C['green']),
        ('06', 'Rejoindre NEBULA', 'Inscription gratuite — formation incluse', C['cyan']),
    ]

    cols = 3
    cw, ch = Inches(3.9), Inches(1.5)
    for i, (num, title, sub, col) in enumerate(items):
        row, col_i = i // cols, i % cols
        x = Inches(0.5) + col_i * (cw + Inches(0.22))
        y = Inches(2.0) + row * (ch + Inches(0.2))

        add_rounded_rect(slide, x, y, cw, ch, C['card2'], radius_pt=10)
        add_rect(slide, x, y, Inches(0.06), ch, col)

        add_text(slide, num, x + Inches(0.2), y + Inches(0.12), Inches(0.7), Inches(0.5),
                 size=28, bold=True, color=col)
        add_text(slide, title, x + Inches(0.2), y + Inches(0.62), cw - Inches(0.3), Inches(0.4),
                 size=14, bold=True, color=C['white'])
        add_text(slide, sub, x + Inches(0.2), y + Inches(0.98), cw - Inches(0.3), Inches(0.44),
                 size=10, color=C['silver'])

    footer(slide, 2, 12, 'Structure des commissions — Confidentiel partenaires')


def slide_03_paliers(prs):
    """Slide 3 — Les 3 paliers"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_dark(slide)

    add_text(slide, '01 — VOS PALIERS DE COMMISSION',
             Inches(0.7), Inches(0.35), Inches(10), Inches(0.5),
             size=12, bold=True, color=C['violet'])
    add_text(slide, 'Plus vous vendez, plus votre taux monte',
             Inches(0.7), Inches(0.8), Inches(11), Inches(0.9),
             size=38, bold=True, color=C['white'])

    paliers = [
        ('STARTER', '25%', '0 – 2 ventes/mois', C['starter'],
         ['Commission immédiate', 'Paiement en 48h', 'Formation gratuite incluse', 'Accès aux outils NEBULA']),
        ('SILVER', '30%', '3 – 6 ventes/mois', C['silver'],
         ['+5 pts par rapport au STARTER', 'Badge partenaire SILVER', 'Support prioritaire WhatsApp', 'Tableau de bord personnel']),
        ('GOLD', '35%', '7+ ventes/mois', C['gold'],
         ['Commission maximale 35%', 'Bonus trimestriel exclusif', 'Accès VIP événements NEBULA', 'Mentorat individuel']),
    ]

    cw = Inches(3.9)
    for i, (name, rate, cond, col, features) in enumerate(paliers):
        x = Inches(0.5) + i * (cw + Inches(0.22))
        y = Inches(1.85)
        ch = Inches(5.2)

        # Carte
        add_rounded_rect(slide, x, y, cw, ch, C['card2'], radius_pt=12)
        # Barre colorée en haut
        add_rounded_rect(slide, x, y, cw, Inches(0.08), col, radius_pt=4)

        # Icône round
        ic_r = Inches(0.45)
        ic_x = x + cw/2 - ic_r
        ic_y = y + Inches(0.3)
        s = slide.shapes.add_shape(9, ic_x, ic_y, ic_r*2, ic_r*2)
        s.fill.solid(); s.fill.fore_color.rgb = col; s.line.fill.background()

        icons = {'STARTER': '★', 'SILVER': '◆', 'GOLD': '♛'}
        add_text(slide, icons[name], ic_x, ic_y + Inches(0.06), ic_r*2, ic_r,
                 size=20, bold=True, color=C['dark'], align=PP_ALIGN.CENTER)

        # Nom palier
        add_text(slide, name, x, y + Inches(1.15), cw, Inches(0.5),
                 size=20, bold=True, color=C['white'], align=PP_ALIGN.CENTER)

        # Taux
        add_text(slide, rate, x, y + Inches(1.62), cw, Inches(0.95),
                 size=56, bold=True, color=col, align=PP_ALIGN.CENTER)

        # Condition
        add_rounded_rect(slide, x + Inches(0.4), y + Inches(2.62),
                         cw - Inches(0.8), Inches(0.34), C['card'], radius_pt=8)
        add_text(slide, cond, x + Inches(0.4), y + Inches(2.64),
                 cw - Inches(0.8), Inches(0.3), size=11, color=col,
                 align=PP_ALIGN.CENTER, bold=True)

        # Séparateur
        add_rect(slide, x + Inches(0.3), y + Inches(3.12),
                 cw - Inches(0.6), Inches(0.015), C['border'])

        # Features
        for j, feat in enumerate(features):
            fy = y + Inches(3.28) + j * Inches(0.43)
            # Tick coloré
            add_text(slide, '✓', x + Inches(0.25), fy, Inches(0.3), Inches(0.36),
                     size=11, bold=True, color=col)
            add_text(slide, feat, x + Inches(0.56), fy, cw - Inches(0.7), Inches(0.36),
                     size=11, color=C['light'])

    footer(slide, 3, 12, 'Structure des commissions — Confidentiel partenaires')


def slide_04_paliers_montants(prs):
    """Slide 4 — Montants concrets par palier sur 150K"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_dark(slide)

    add_text(slide, '01 — EN CHIFFRES RÉELS',
             Inches(0.7), Inches(0.35), Inches(10), Inches(0.5),
             size=12, bold=True, color=C['violet'])
    add_text(slide, 'Sur une vitrine à 150 000 FCFA',
             Inches(0.7), Inches(0.8), Inches(10), Inches(0.9),
             size=38, bold=True, color=C['white'])

    paliers = [
        ('STARTER', '25%', '37 500', C['starter']),
        ('SILVER',  '30%', '45 000', C['silver']),
        ('GOLD',    '35%', '52 500', C['gold']),
    ]

    bw = Inches(3.6)
    bh = Inches(2.5)
    for i, (name, rate, amount, col) in enumerate(paliers):
        x = Inches(0.8) + i * (bw + Inches(0.47))
        y = Inches(2.1)
        add_rounded_rect(slide, x, y, bw, bh, C['card2'], radius_pt=12)
        add_rect(slide, x, y, bw, Inches(0.06), col)

        add_text(slide, name, x, y + Inches(0.2), bw, Inches(0.46),
                 size=18, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_text(slide, rate, x, y + Inches(0.65), bw, Inches(0.8),
                 size=42, bold=True, color=C['white'], align=PP_ALIGN.CENTER)
        add_text(slide, 'sur vos ventes directes', x, y + Inches(1.4), bw, Inches(0.35),
                 size=10, color=C['silver'], align=PP_ALIGN.CENTER)

        # Montant
        add_rounded_rect(slide, x + Inches(0.3), y + Inches(1.82),
                         bw - Inches(0.6), Inches(0.54), col, radius_pt=8)
        add_text(slide, f'{amount} FCFA',
                 x + Inches(0.3), y + Inches(1.9), bw - Inches(0.6), Inches(0.38),
                 size=18, bold=True, color=C['dark'], align=PP_ALIGN.CENTER)

    # Flèche progression
    add_text(slide, '→', Inches(4.4), Inches(3.0), Inches(0.6), Inches(0.6),
             size=32, bold=True, color=C['violet'], align=PP_ALIGN.CENTER)
    add_text(slide, '→', Inches(8.0), Inches(3.0), Inches(0.6), Inches(0.6),
             size=32, bold=True, color=C['violet'], align=PP_ALIGN.CENTER)

    # Bloc info bas
    info_data = [
        ('Paiement garanti', 'Sous 48h après validation de la vente', C['green']),
        ('Commission directe', 'Créditée sur votre compte dès la confirmation', C['cyan']),
        ('Palier recalculé', 'Chaque 1er du mois selon vos ventes du mois précédent', C['violet']),
    ]
    for i, (title, sub, col) in enumerate(info_data):
        x = Inches(0.8) + i * Inches(4.24)
        y = Inches(5.1)
        add_rounded_rect(slide, x, y, Inches(3.9), Inches(1.8), C['card'], radius_pt=10)
        add_rect(slide, x, y, Inches(0.06), Inches(1.8), col)
        add_text(slide, title, x + Inches(0.2), y + Inches(0.2), Inches(3.5), Inches(0.42),
                 size=14, bold=True, color=C['white'])
        add_text(slide, sub, x + Inches(0.2), y + Inches(0.65), Inches(3.5), Inches(0.9),
                 size=11, color=C['silver'])

    footer(slide, 4, 12, 'Structure des commissions — Confidentiel partenaires')


def slide_05_reseau(prs):
    """Slide 5 — Réseau 3 niveaux concept"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_dark(slide)

    add_text(slide, '02 — REVENUS DE RÉSEAU',
             Inches(0.7), Inches(0.35), Inches(10), Inches(0.5),
             size=12, bold=True, color=C['violet'])
    add_text(slide, 'Gagnez même quand vous ne vendez pas',
             Inches(0.7), Inches(0.8), Inches(11), Inches(0.9),
             size=38, bold=True, color=C['white'])

    # Concept central
    add_text(slide, 'Chaque personne que vous recrutez dans le réseau NEBULA\nvous génère des commissions sur TOUTES ses ventes — automatiquement.',
             Inches(0.7), Inches(1.78), Inches(11.5), Inches(0.9),
             size=15, color=C['silver'])

    levels = [
        ('N1', '5%', 'Ventes de vos filleuls directs',
         'Personnes que vous avez recrutées', C['violet'], 7500, Inches(0.8)),
        ('N2', '3%', 'Ventes des filleuls de vos filleuls',
         'Niveau 2 de profondeur', C['starter'], 4500, Inches(4.6)),
        ('N3', '2%', 'Ventes au 3ème niveau de profondeur',
         'Le réseau travaille pour vous', C['green'], 3000, Inches(8.4)),
    ]

    for level, rate, title, sub, col, amount, x in levels:
        y = Inches(2.85)
        bw, bh = Inches(4.0), Inches(3.7)
        add_rounded_rect(slide, x, y, bw, bh, C['card2'], radius_pt=12)

        # Badge niveau
        add_rounded_rect(slide, x + Inches(0.3), y + Inches(0.25),
                         Inches(0.7), Inches(0.38), col, radius_pt=8)
        add_text(slide, level, x + Inches(0.3), y + Inches(0.26), Inches(0.7), Inches(0.36),
                 size=16, bold=True, color=C['dark'], align=PP_ALIGN.CENTER)

        # Taux
        add_text(slide, rate, x + Inches(1.15), y + Inches(0.27), Inches(2.6), Inches(0.4),
                 size=26, bold=True, color=col)

        # Titre
        add_text(slide, title, x + Inches(0.25), y + Inches(0.85), bw - Inches(0.4), Inches(0.55),
                 size=14, bold=True, color=C['white'])

        add_text(slide, sub, x + Inches(0.25), y + Inches(1.35), bw - Inches(0.4), Inches(0.5),
                 size=11, color=C['silver'])

        # Montant concret
        add_rect(slide, x + Inches(0.25), y + Inches(2.05),
                 bw - Inches(0.5), Inches(0.015), C['border'])
        add_text(slide, 'Sur une vente de 150 000 F :',
                 x + Inches(0.25), y + Inches(2.2), bw - Inches(0.4), Inches(0.36),
                 size=10, color=C['muted'])
        add_text(slide, f'+ {amount:,} FCFA'.replace(',', ' '),
                 x + Inches(0.25), y + Inches(2.55), bw - Inches(0.4), Inches(0.6),
                 size=22, bold=True, color=col)
        add_text(slide, 'pour vous, automatiquement',
                 x + Inches(0.25), y + Inches(3.12), bw - Inches(0.4), Inches(0.35),
                 size=10, color=C['silver'])

    # Connecteur visuel entre les niveaux
    for arr_x in [Inches(4.85), Inches(8.65)]:
        add_text(slide, '+', arr_x, Inches(4.3), Inches(0.5), Inches(0.5),
                 size=28, bold=True, color=C['border'], align=PP_ALIGN.CENTER)

    footer(slide, 5, 12, 'Structure des commissions — Confidentiel partenaires')


def slide_06_arbre_shad(prs):
    """Slide 6 — Arbre réseau Shad"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_dark(slide)

    add_text(slide, "03 — L'ARBRE RÉSEAU DE SHAD",
             Inches(0.7), Inches(0.35), Inches(10), Inches(0.5),
             size=12, bold=True, color=C['violet'])
    add_text(slide, "Visualisez comment les commissions circulent",
             Inches(0.7), Inches(0.8), Inches(11), Inches(0.9),
             size=34, bold=True, color=C['white'])

    # Fond carte centrale
    add_rounded_rect(slide, Inches(0.4), Inches(1.75), Inches(12.5), Inches(5.45),
                     C['card2'], radius_pt=14)

    # Positions des nœuds (cx, cy en Inches dans le slide)
    nodes = {
        'marc':  (6.65, 2.55,  0.42, C['violet'], 'MARC',  'Votre recruteur', 'MA', False, ''),
        'shad':  (6.65, 4.05,  0.52, C['gold'],   'SHAD',  'SILVER — 30%',    'SH', True,  'VOUS'),
        'paul':  (4.2,  5.65,  0.4,  C['starter'],'PAUL',  'STARTER — 25%',   'PA', False, ''),
        'mia':   (9.1,  5.65,  0.4,  C['pink'],   'MIA',   'STARTER — 25%',   'MI', False, ''),
        'kofi':  (2.8,  7.0,   0.33, C['green'],  'KOFI',  'STARTER — 25%',   'KO', False, ''),
    }

    # Connecteurs (lignes)
    lines = [
        ('marc', 'shad', C['violet'], 'N1 → 5%'),
        ('shad', 'paul', C['starter'], 'N1 → 5%'),
        ('shad', 'mia',  C['pink'],   'N1 → 5%'),
        ('paul', 'kofi', C['green'],  'N2 → 3%'),
    ]

    for src, dst, col, lbl in lines:
        sx, sy = nodes[src][0], nodes[src][1]
        dx, dy = nodes[dst][0], nodes[dst][1]
        add_connector_line(slide, sx, sy + nodes[src][2] + 0.02,
                           dx, dy - nodes[dst][2] - 0.02, col, 2)
        mx = (sx + dx) / 2
        my = (sy + nodes[src][2] + 0.02 + dy - nodes[dst][2] - 0.02) / 2
        add_label_pill_on_line(slide, mx, my, lbl, col)

    # Nœuds
    for key, (cx, cy, r, col, name, role, init, hl, badge) in nodes.items():
        add_circle_node(slide, cx, cy, r, col, name, role,
                        size_top=12, size_bot=10,
                        initials=init, highlight=hl, badge_text=badge)

    # Légende droite
    leg_x = Inches(11.5)
    add_text(slide, 'LÉGENDE', leg_x, Inches(2.0), Inches(1.5), Inches(0.4),
             size=10, bold=True, color=C['muted'])
    leg_items = [
        (C['violet'], 'Recruteur N0'),
        (C['gold'],   'Vous (SILVER)'),
        (C['starter'],'Filleul direct N1'),
        (C['green'],  'Filleul N2'),
    ]
    for i, (col, label) in enumerate(leg_items):
        ly = Inches(2.5) + i * Inches(0.6)
        s = slide.shapes.add_shape(9, leg_x, ly, Inches(0.22), Inches(0.22))
        s.fill.solid(); s.fill.fore_color.rgb = col; s.line.fill.background()
        add_text(slide, label, leg_x + Inches(0.3), ly - Inches(0.02),
                 Inches(1.2), Inches(0.3), size=9, color=C['light'])

    footer(slide, 6, 12, 'Structure des commissions — Confidentiel partenaires')


def slide_07_gains_shad(prs):
    """Slide 7 — Gains de Shad"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_dark(slide)

    add_text(slide, "03 — CE QUE SHAD GAGNE CE MOIS",
             Inches(0.7), Inches(0.35), Inches(10), Inches(0.5),
             size=12, bold=True, color=C['gold'])
    add_text(slide, "5 ventes directes + réseau actif",
             Inches(0.7), Inches(0.8), Inches(10), Inches(0.9),
             size=38, bold=True, color=C['white'])

    rows = [
        ('Ses 5 ventes directes',  'SILVER 30%', '5 × 45 000 F', '225 000 FCFA', C['gold'],    1.0),
        ('Ventes de Paul',         'N1 — 5%',    '3 × 7 500 F',  '22 500 FCFA',  C['starter'], 0.10),
        ('Ventes de Mia',          'N1 — 5%',    '2 × 7 500 F',  '15 000 FCFA',  C['pink'],    0.067),
        ('Ventes de Kofi',         'N2 — 3%',    '2 × 4 500 F',  '9 000 FCFA',   C['green'],   0.04),
    ]

    max_w = Inches(7.0)
    row_h = Inches(0.95)
    for i, (source, badge_txt, calcul, montant, col, pct) in enumerate(rows):
        y = Inches(2.05) + i * (row_h + Inches(0.12))
        add_rounded_rect(slide, Inches(0.7), y, Inches(12.0), row_h, C['card2'], radius_pt=8)
        add_rect(slide, Inches(0.7), y, Inches(0.06), row_h, col)

        # Source
        add_text(slide, source, Inches(1.0), y + Inches(0.18), Inches(3.8), Inches(0.46),
                 size=15, bold=True, color=C['white'])

        # Badge
        add_rounded_rect(slide, Inches(1.0), y + Inches(0.6), Inches(1.3), Inches(0.26),
                         col, radius_pt=6)
        add_text(slide, badge_txt, Inches(1.0), y + Inches(0.61), Inches(1.3), Inches(0.24),
                 size=9, bold=True, color=C['dark'], align=PP_ALIGN.CENTER)

        # Calcul
        add_text(slide, calcul, Inches(5.2), y + Inches(0.28), Inches(2.5), Inches(0.4),
                 size=13, color=C['silver'], align=PP_ALIGN.CENTER)

        # Barre proportionnelle
        bar_x = Inches(5.2)
        bar_y = y + Inches(0.62)
        bar_max_w = Inches(3.8)
        bar_h_px = Inches(0.18)
        add_rounded_rect(slide, bar_x, bar_y, bar_max_w, bar_h_px, C['card'], radius_pt=4)
        add_rounded_rect(slide, bar_x, bar_y, bar_max_w * pct, bar_h_px, col, radius_pt=4)

        # Montant
        add_text(slide, montant, Inches(9.2), y + Inches(0.22), Inches(3.2), Inches(0.52),
                 size=20, bold=True, color=col, align=PP_ALIGN.RIGHT)

    # Total
    add_rounded_rect(slide, Inches(0.7), Inches(6.18), Inches(12.0), Inches(0.95),
                     RGBColor(0x3D, 0x2C, 0x00), radius_pt=10)
    add_rect(slide, Inches(0.7), Inches(6.18), Inches(12.0), Inches(0.05), C['gold'])
    add_text(slide, 'TOTAL SHAD CE MOIS', Inches(1.0), Inches(6.3), Inches(7), Inches(0.5),
             size=18, bold=True, color=C['white'])
    add_text(slide, '271 500 FCFA', Inches(9.2), Inches(6.3), Inches(3.2), Inches(0.5),
             size=28, bold=True, color=C['gold'], align=PP_ALIGN.RIGHT)

    footer(slide, 7, 12, 'Structure des commissions — Confidentiel partenaires')


def slide_08_gains_marc(prs):
    """Slide 8 — Ce que Marc gagne (revenu passif)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_dark(slide)

    add_text(slide, "04 — LE REVENU PASSIF DE MARC",
             Inches(0.7), Inches(0.35), Inches(10), Inches(0.5),
             size=12, bold=True, color=C['violet'])
    add_text(slide, "Marc n'a rien vendu ce mois — et il encaisse quand même",
             Inches(0.7), Inches(0.8), Inches(11), Inches(0.9),
             size=32, bold=True, color=C['white'])

    # Bloc intro
    add_rounded_rect(slide, Inches(0.7), Inches(1.82), Inches(12.0), Inches(0.65),
                     C['card'], radius_pt=8)
    add_text(slide,
             "Marc a recruté Shad. Shad a recruté Paul et Mia. Paul a recruté Kofi. Marc touche des commissions sur toute la chaîne.",
             Inches(1.0), Inches(1.9), Inches(11.4), Inches(0.5),
             size=12, color=C['light'])

    rows = [
        ('Ventes de Shad',  'N1 — 5%', '5 × 7 500 F', '37 500 FCFA', C['gold']),
        ('Ventes de Paul',  'N2 — 3%', '3 × 4 500 F', '13 500 FCFA', C['starter']),
        ('Ventes de Mia',   'N2 — 3%', '2 × 4 500 F', '9 000 FCFA',  C['pink']),
        ('Ventes de Kofi',  'N3 — 2%', '2 × 3 000 F', '6 000 FCFA',  C['green']),
    ]

    row_h = Inches(0.88)
    for i, (source, badge_txt, calcul, montant, col) in enumerate(rows):
        y = Inches(2.65) + i * (row_h + Inches(0.1))
        add_rounded_rect(slide, Inches(0.7), y, Inches(12.0), row_h, C['card2'], radius_pt=8)
        add_rect(slide, Inches(0.7), y, Inches(0.06), row_h, col)

        add_text(slide, source, Inches(1.0), y + Inches(0.14), Inches(3.8), Inches(0.42),
                 size=15, bold=True, color=C['white'])

        add_rounded_rect(slide, Inches(1.0), y + Inches(0.55), Inches(1.3), Inches(0.24),
                         col, radius_pt=6)
        add_text(slide, badge_txt, Inches(1.0), y + Inches(0.555), Inches(1.3), Inches(0.23),
                 size=9, bold=True, color=C['dark'], align=PP_ALIGN.CENTER)

        add_text(slide, calcul, Inches(5.5), y + Inches(0.24), Inches(3.0), Inches(0.38),
                 size=13, color=C['silver'], align=PP_ALIGN.CENTER)

        add_text(slide, montant, Inches(9.2), y + Inches(0.19), Inches(3.2), Inches(0.48),
                 size=20, bold=True, color=col, align=PP_ALIGN.RIGHT)

    # Total + comparaison
    add_rounded_rect(slide, Inches(0.7), Inches(6.28), Inches(5.8), Inches(0.9),
                     RGBColor(0x1E, 0x1B, 0x4B), radius_pt=10)
    add_rect(slide, Inches(0.7), Inches(6.28), Inches(5.8), Inches(0.04), C['violet'])
    add_text(slide, 'TOTAL MARC (revenu passif)', Inches(1.0), Inches(6.38), Inches(3.5), Inches(0.42),
             size=13, bold=True, color=C['silver'])
    add_text(slide, '66 000 FCFA', Inches(1.0), Inches(6.68), Inches(5.1), Inches(0.38),
             size=26, bold=True, color=C['violet'], align=PP_ALIGN.RIGHT)

    add_rounded_rect(slide, Inches(6.8), Inches(6.28), Inches(5.9), Inches(0.9),
                     RGBColor(0x0F, 0x0B, 0x24), radius_pt=10)
    add_text(slide, '💡  Sans faire une seule vente ce mois',
             Inches(7.1), Inches(6.38), Inches(5.5), Inches(0.42),
             size=13, bold=True, color=C['cyan'])
    add_text(slide, 'C\'est le pouvoir des revenus résiduels de réseau.',
             Inches(7.1), Inches(6.7), Inches(5.5), Inches(0.38),
             size=11, color=C['silver'])

    footer(slide, 8, 12, 'Structure des commissions — Confidentiel partenaires')


def slide_09_flux_vente(prs):
    """Slide 9 — Flux d'une vente (répartition)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_dark(slide)

    add_text(slide, "05 — RÉPARTITION SUR UNE VENTE",
             Inches(0.7), Inches(0.35), Inches(10), Inches(0.5),
             size=12, bold=True, color=C['pink'])
    add_text(slide, "Chaque vente de 150 000 FCFA est répartie ainsi",
             Inches(0.7), Inches(0.8), Inches(11), Inches(0.9),
             size=36, bold=True, color=C['white'])

    # Segments visuels (horizontal stacked bar)
    segments = [
        ('Vendeur GOLD', '35%', '52 500 F', C['gold'],    0.35),
        ('N1 recruteur', '5%',  '7 500 F',  C['violet'],  0.05),
        ('N2',           '3%',  '4 500 F',  C['starter'], 0.03),
        ('N3',           '2%',  '3 000 F',  C['green'],   0.02),
        ('NEBULA Agency','55%', '82 500 F', C['pink'],    0.55),
    ]

    bar_x = Inches(0.7)
    bar_y = Inches(2.2)
    bar_w = Inches(11.9)
    bar_h = Inches(1.2)
    add_rounded_rect(slide, bar_x, bar_y, bar_w, bar_h, C['card2'], radius_pt=10)

    cur_x = bar_x
    for name, rate, amount, col, pct in segments:
        seg_w = bar_w * pct
        add_rounded_rect(slide, cur_x, bar_y, seg_w, bar_h, col, radius_pt=8)
        if pct > 0.06:
            add_text(slide, rate, cur_x, bar_y + Inches(0.35), seg_w, Inches(0.5),
                     size=18, bold=True, color=C['dark'], align=PP_ALIGN.CENTER)
        cur_x += seg_w

    # Légendes sous la barre
    cur_x = bar_x
    for name, rate, amount, col, pct in segments:
        seg_w = bar_w * pct
        lx = cur_x + seg_w / 2
        ly = bar_y + bar_h + Inches(0.18)
        # Trait vertical
        add_rect(slide, lx - Inches(0.01), bar_y + bar_h, Inches(0.02), Inches(0.18), col)
        lbl_w = max(seg_w, Inches(0.9))
        add_text(slide, name, lx - lbl_w/2, ly, lbl_w, Inches(0.35),
                 size=9, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_text(slide, amount, lx - lbl_w/2, ly + Inches(0.35), lbl_w, Inches(0.32),
                 size=11, bold=True, color=C['white'], align=PP_ALIGN.CENTER)
        cur_x += seg_w

    # Cards récap bas
    recap = [
        ('Scénario : GOLD seul', 'Vendeur + réseau sur 2 niveaux', '35% sortis', '97 500 F gardés', C['gold']),
        ('Scénario : SILVER + réseau', '30% + 5% + 3% = 38% sortis', '38% sortis', '93 000 F gardés', C['silver']),
        ('Scénario : STARTER + réseau complet', '25% + 5% + 3% + 2% = 35%', '35% sortis', '97 500 F gardés', C['starter']),
    ]
    cw = Inches(4.1)
    for i, (title, sub, out, kept, col) in enumerate(recap):
        x = Inches(0.5) + i * (cw + Inches(0.34))
        y = Inches(5.05)
        add_rounded_rect(slide, x, y, cw, Inches(2.15), C['card2'], radius_pt=10)
        add_rect(slide, x, y, cw, Inches(0.05), col)
        add_text(slide, title, x + Inches(0.2), y + Inches(0.18), cw - Inches(0.3), Inches(0.42),
                 size=13, bold=True, color=C['white'])
        add_text(slide, sub, x + Inches(0.2), y + Inches(0.62), cw - Inches(0.3), Inches(0.42),
                 size=10, color=C['silver'])
        add_text(slide, kept, x + Inches(0.2), y + Inches(1.35), cw - Inches(0.3), Inches(0.55),
                 size=18, bold=True, color=C['pink'])
        add_text(slide, 'conservés par NEBULA', x + Inches(0.2), y + Inches(1.82),
                 cw - Inches(0.3), Inches(0.3), size=9, color=C['muted'])

    footer(slide, 9, 12, 'Structure des commissions — Confidentiel partenaires')


def slide_10_simulation(prs):
    """Slide 10 — Simulation revenus"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_dark(slide)

    add_text(slide, "06 — SIMULEZ VOS REVENUS",
             Inches(0.7), Inches(0.35), Inches(10), Inches(0.5),
             size=12, bold=True, color=C['green'])
    add_text(slide, "Où pouvez-vous être dans 3 mois ?",
             Inches(0.7), Inches(0.8), Inches(10), Inches(0.9),
             size=38, bold=True, color=C['white'])

    scenarios = [
        ('Débutant',   '2 ventes', 'STARTER', '60 000 F',  '+ 15 000 F',  '75 000 FCFA',   C['muted'],   0.076,  '★'),
        ('Actif',      '5 ventes', 'SILVER',  '150 000 F', '+ 45 000 F',  '195 000 FCFA',  C['starter'], 0.197,  '★★'),
        ('Champion',   '10 ventes','SILVER',  '300 000 F', '+ 90 000 F',  '390 000 FCFA',  C['violet'],  0.394,  '★★★'),
        ('GOLD Élite', '15 ventes','GOLD',    '787 500 F', '+ 200 000 F', '987 500 FCFA',  C['gold'],    1.0,    '♛'),
    ]

    bar_max = Inches(4.5)
    row_h = Inches(1.28)
    for i, (name, ventes, palier, direct, reseau, total, col, pct, stars) in enumerate(scenarios):
        y = Inches(2.0) + i * (row_h + Inches(0.1))

        # Fond ligne
        add_rounded_rect(slide, Inches(0.5), y, Inches(12.4), row_h, C['card2'], radius_pt=10)
        if col == C['gold']:
            add_rect(slide, Inches(0.5), y, Inches(12.4), Inches(0.05), col)

        # Stars / nom
        add_text(slide, stars, Inches(0.75), y + Inches(0.12), Inches(1.2), Inches(0.36),
                 size=14, color=col)
        add_text(slide, name, Inches(0.75), y + Inches(0.48), Inches(2.0), Inches(0.48),
                 size=16, bold=True, color=C['white'])

        # Badge palier
        add_rounded_rect(slide, Inches(0.75), y + Inches(0.9), Inches(1.3), Inches(0.26), col, radius_pt=6)
        add_text(slide, palier, Inches(0.75), y + Inches(0.905), Inches(1.3), Inches(0.24),
                 size=9, bold=True, color=C['dark'], align=PP_ALIGN.CENTER)

        # Ventes
        add_text(slide, ventes, Inches(2.9), y + Inches(0.42), Inches(1.5), Inches(0.42),
                 size=14, bold=True, color=C['silver'], align=PP_ALIGN.CENTER)

        # Direct + réseau
        add_text(slide, direct, Inches(4.5), y + Inches(0.2), Inches(2.0), Inches(0.4),
                 size=14, bold=True, color=C['light'], align=PP_ALIGN.CENTER)
        add_text(slide, reseau, Inches(4.5), y + Inches(0.65), Inches(2.0), Inches(0.38),
                 size=12, color=C['cyan'], align=PP_ALIGN.CENTER)

        # Barre
        bar_x = Inches(6.8)
        bar_y_c = y + row_h / 2 - Inches(0.1)
        add_rounded_rect(slide, bar_x, bar_y_c, bar_max, Inches(0.22), C['card'], radius_pt=6)
        add_rounded_rect(slide, bar_x, bar_y_c, bar_max * pct, Inches(0.22), col, radius_pt=6)

        # Total
        add_text(slide, total, Inches(11.5), y + Inches(0.34), Inches(1.3), Inches(0.56),
                 size=18, bold=True, color=col, align=PP_ALIGN.RIGHT)

    # En-têtes colonnes
    headers = [('PROFIL', Inches(0.75)), ('VENTES', Inches(2.9)), ('REVENUS', Inches(4.5)),
               ('PROGRESSION', Inches(6.8)), ('TOTAL', Inches(11.5))]
    for hdr, hx in headers:
        add_text(slide, hdr, hx, Inches(1.7), Inches(2.0), Inches(0.34),
                 size=9, bold=True, color=C['muted'])

    footer(slide, 10, 12, 'Structure des commissions — Confidentiel partenaires')


def slide_11_avantages(prs):
    """Slide 11 — Pourquoi rejoindre NEBULA"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_dark(slide)

    add_text(slide, "07 — POURQUOI NEBULA ?",
             Inches(0.7), Inches(0.35), Inches(10), Inches(0.5),
             size=12, bold=True, color=C['cyan'])
    add_text(slide, "Un produit qui se vend, un système qui paie",
             Inches(0.7), Inches(0.8), Inches(11), Inches(0.9),
             size=38, bold=True, color=C['white'])

    avantages = [
        ('Produit concret', 'Vous vendez des vraies vitrines web à de vraies PME. Pas un abonnement virtuel.', '🌐', C['cyan']),
        ('Marché immense', 'Des milliers de PME en Afrique de l\'Ouest sans présence digitale. Demande réelle.', '📈', C['green']),
        ('Commission dès J+1', 'Vente validée le lundi → argent sur votre compte mercredi. Aucune attente.', '⚡', C['gold']),
        ('Formation incluse', 'On vous forme à vendre, à pitcher, à suivre vos clients. Zéro expérience requise.', '🎓', C['violet']),
        ('Réseau = rente', 'Chaque partenaire que vous recrutez grossit votre revenu mensuel automatiquement.', '♾️', C['pink']),
        ('Support WhatsApp', 'Une équipe NEBULA dédiée répond 7j/7 à toutes vos questions de vente.', '💬', C['starter']),
    ]

    cw = Inches(3.9)
    ch = Inches(1.9)
    for i, (title, desc, icon, col) in enumerate(avantages):
        row, col_i = i // 3, i % 3
        x = Inches(0.5) + col_i * (cw + Inches(0.22))
        y = Inches(2.05) + row * (ch + Inches(0.2))

        add_rounded_rect(slide, x, y, cw, ch, C['card2'], radius_pt=10)
        add_rect(slide, x, y, Inches(0.05), ch, col)

        add_text(slide, icon, x + Inches(0.3), y + Inches(0.15), Inches(0.55), Inches(0.55),
                 size=22, align=PP_ALIGN.CENTER)
        add_text(slide, title, x + Inches(0.3), y + Inches(0.62), cw - Inches(0.45), Inches(0.42),
                 size=14, bold=True, color=C['white'])
        add_text(slide, desc, x + Inches(0.3), y + Inches(1.07), cw - Inches(0.45), Inches(0.76),
                 size=10, color=C['silver'])

    footer(slide, 11, 12, 'Structure des commissions — Confidentiel partenaires')


def slide_12_cta(prs):
    """Slide 12 — CTA final"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg_dark(slide)

    # Grand fond dégradé simulé
    add_rounded_rect(slide, Inches(0.5), Inches(0.4), Inches(12.3), Inches(6.65),
                     C['card'], radius_pt=18)
    add_rect(slide, Inches(0.5), Inches(0.4), Inches(12.3), Inches(0.06), C['purple'])
    add_rect(slide, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.06), C['pink'])

    # Cercles déco
    for (cx, cy, cr, col) in [
        (11.5, 1.2, 1.4, RGBColor(0x2D, 0x1A, 0x6E)),
        (1.5, 6.5, 1.0, RGBColor(0x3D, 0x18, 0x35)),
        (10.5, 6.8, 0.8, RGBColor(0x1A, 0x3A, 0x2E)),
    ]:
        s = slide.shapes.add_shape(9, Inches(cx - cr), Inches(cy - cr),
                                   Inches(cr*2), Inches(cr*2))
        s.fill.solid(); s.fill.fore_color.rgb = col; s.line.fill.background()

    # Badge
    add_badge(slide, '🚀  REJOIGNEZ NEBULA AGENCY', Inches(4.5), Inches(1.1), C['purple'])

    # Titre
    add_text(slide, 'Commencez à gagner', Inches(1.5), Inches(1.75),
             Inches(10), Inches(1.1), size=52, bold=True, color=C['white'],
             align=PP_ALIGN.CENTER)
    add_text(slide, 'dès votre 1ère vente', Inches(1.5), Inches(2.75),
             Inches(10), Inches(1.0), size=52, bold=True, color=C['violet'],
             align=PP_ALIGN.CENTER)

    # Sous-titre
    add_text(slide,
             'Inscription gratuite  •  Formation incluse  •  Paiement sous 48h  •  Réseau illimité',
             Inches(1.5), Inches(3.85), Inches(10), Inches(0.5),
             size=14, color=C['silver'], align=PP_ALIGN.CENTER)

    # 3 boutons d'action
    actions = [
        ('WhatsApp', '+229 97 XX XX XX', C['green']),
        ('Email', 'contact@nebula-agency.com', C['cyan']),
        ('Inscription', 'nebula-agency.com/partenaire', C['violet']),
    ]
    bw = Inches(3.5)
    for i, (label, contact, col) in enumerate(actions):
        x = Inches(1.4) + i * (bw + Inches(0.5))
        y = Inches(4.65)
        add_rounded_rect(slide, x, y, bw, Inches(1.4), C['card2'], radius_pt=12)
        add_rect(slide, x, y, bw, Inches(0.05), col)
        add_text(slide, label, x, y + Inches(0.15), bw, Inches(0.46),
                 size=13, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_text(slide, contact, x, y + Inches(0.62), bw, Inches(0.5),
                 size=12, color=C['light'], align=PP_ALIGN.CENTER)

    # Quote bas
    add_text(slide,
             '"Le meilleur moment pour planter un arbre, c\'était hier. Le second meilleur moment, c\'est maintenant."',
             Inches(2.0), Inches(6.2), Inches(9.3), Inches(0.6),
             size=11, italic=True, color=C['muted'], align=PP_ALIGN.CENTER)

    footer(slide, 12, 12, 'NEBULA Agency — Cotonou, Bénin — Afrique de l\'Ouest')


# ─── Main ─────────────────────────────────────────────────────────

def build_pptx(output_path):
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    print('Slide 1  — Couverture...')
    slide_01_cover(prs)
    print('Slide 2  — Agenda...')
    slide_02_agenda(prs)
    print('Slide 3  — Les 3 paliers...')
    slide_03_paliers(prs)
    print('Slide 4  — Montants concrets...')
    slide_04_paliers_montants(prs)
    print('Slide 5  — Réseau 3 niveaux...')
    slide_05_reseau(prs)
    print('Slide 6  — Arbre Shad...')
    slide_06_arbre_shad(prs)
    print('Slide 7  — Gains Shad...')
    slide_07_gains_shad(prs)
    print('Slide 8  — Gains Marc (passif)...')
    slide_08_gains_marc(prs)
    print('Slide 9  — Flux d\'une vente...')
    slide_09_flux_vente(prs)
    print('Slide 10 — Simulation revenus...')
    slide_10_simulation(prs)
    print('Slide 11 — Pourquoi NEBULA...')
    slide_11_avantages(prs)
    print('Slide 12 — CTA final...')
    slide_12_cta(prs)

    prs.save(output_path)
    print(f'\nPPTX généré : {output_path}')


if __name__ == '__main__':
    output = '/home/user/nebula-agency/_docs/NEBULA_Commission_Conference_2025.pptx'
    build_pptx(output)
